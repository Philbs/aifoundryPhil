from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

# Dados Detalhados M2 (Hotel Agent)
M02 = [
    {"title": "Building Your Hotel Agent", "description": "Módulo 2: Do conceito à criação. Aprenda a configurar um agente profissional.", "image": "img24.png"},
    {"title": "Lições de Sucesso Enterprise", "description": "Fatores chave: Iteração, Dados Reais, Confiança e Foco em Escalabilidade.", "image": "img15.png"},
    {"title": "Caso de Estudo: Tax Assistant", "description": "Um modelo para o seu Hotel Agent: começar simples e iterar com feedback.", "image": "img10.png"},
    {"title": "Deploying the Model", "description": "Escolha GPT-4 Balanced para um equilíbrio perfeito entre inteligência e custo.", "image": "img22.png"},
    {"title": "O Setup Panel", "description": "Seu centro de comando: defina nome, implantação e comece a moldar a IA.", "image": "img13.png"},
    {"title": "System Instructions", "description": "Defina o coração do agente: 'Você é um concierge simpático'. Regras e tons claros.", "image": "img14.png"},
    {"title": "Parâmetro: Temperature", "description": "0.0 é factual; 2.0 é criativo. Para hotéis, use 0.2 para fatos e 0.5 para equilíbrio.", "image": "img11.png"},
    {"title": "Parâmetro: Top P", "description": "Nucleus Sampling: limita a escolha do modelo às palavras mais prováveis.", "image": "img20.png"},
    {"title": "Grounding (Aterramento)", "description": "Conecte seu agente a documentos oficiais do hotel para respostas reais.", "image": "img19.png"},
    {"title": "Aba Evaluation", "description": "O laboratório de comparação: teste versões Formal vs Casual lado a lado.", "image": "img31.png"}
]

# Dados Detalhados M3 (AI Services)
M03 = [
    {"title": "AI Services Enhancement", "description": "Módulo 3: Dê sentidos aos seus agentes. Explore NLU, visão e fala.", "image": "img27.png"},
    {"title": "O Poder do NLU", "description": "Entenda intenção e contexto, não apenas palavras-chave isoladas.", "image": "img26.png"},
    {"title": "Intent Recognition", "description": "Identifique o 'porquê' do hóspede. Reduza erros de roteamento em até 60%.", "image": "img5.png"},
    {"title": "Summarization Call Center", "description": "Mapeie automaticamente: Recapitulativo, Problema e Resolução final.", "image": "img32.png"},
    {"title": "Named Entity Recognition", "description": "Extração automática de nomes, datas e locais diretamente do texto.", "image": "img26.png"},
    {"title": "Dynamic Language Detection", "description": "Detecte idiomas e traduza em tempo real para um alcance global sem barreiras.", "image": "img5.png"},
    {"title": "Azure AI Vision", "description": "Permita que o agente enxergue passaportes e documentos via OCR avançado.", "image": "img12.png"},
    {"title": "Azure AI Speech", "description": "Sintetize vozes humanas e empáticas para quiosques e suporte por voz.", "image": "img33.png"},
    {"title": "Contextual Grounding", "description": "Integre visão e fala com o cérebro do agente para uma percepção 360 graus.", "image": "img24.png"}
]

# Dados Detalhados M4 (Professional Dev)
M04 = [
    {"title": "Professional Development", "description": "Módulo 4: Qualidade e Segurança. Teste e lance agentes em nível mundial.", "image": "img23.png"},
    {"title": "Impacto das Falhas de IA", "description": "Erros em bots escalam para desastres financeiros, judiciais e de reputação.", "image": "img10.png"},
    {"title": "Caso Real: Companhia Aérea", "description": "Política de reembolso errada gerou processo judicial e perda de confiança.", "image": "img28.png"},
    {"title": "Caso Real: Chatbot Bancário", "description": "Vazamento de dados privados gerou multas pesadas e crise da marca.", "image": "img29.png"},
    {"title": "Functional Testing", "description": "Valide respostas sobre check-out e Wi-Fi de forma sistemática.", "image": "img1.png"},
    {"title": "Adversarial Testing", "description": "Tente 'quebrar' o agente com injeção de prompts para testar segurança.", "image": "img30.png"},
    {"title": "Prompt Injection", "description": "Ameaça real: usuários tentam assumir controle do bot via prompts maliciosos.", "image": "img30.png"},
    {"title": "Deployment: Canary Deploy", "description": "Lance para 5% dos usuários primeiro para validar antes do rollout total.", "image": "img22.png"},
    {"title": "Responsible AI Standards", "description": "Princípios: Justiça, Confiabilidade, Privacidade, Inclusão e Transparência.", "image": "img27.png"}
]

# M1 (Resumido para o script)
M01 = [
    {"title": "Azure AI Foundry", "description": "Introdução ao ecossistema de agentes inteligentes da Microsoft.", "image": "img14.png"},
    {"title": "Arquitetura Unificada", "description": "Escalabilidade global e segurança robusta para o ciclo de vida da IA.", "image": "img22.png"}
]

def create_pptx(module_name, slides_data):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    images_path = os.path.join("public", "images")

    for data in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(11, 14, 20)
        
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(1.5))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = data["title"]
        p.font.bold = True
        p.font.size = Pt(40)
        p.font.color.rgb = RGBColor(0, 120, 212)
        
        txBoxDesc = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(6), Inches(4))
        tfDesc = txBoxDesc.text_frame
        tfDesc.word_wrap = True
        pDesc = tfDesc.add_paragraph()
        pDesc.text = data["description"]
        pDesc.font.size = Pt(24)
        pDesc.font.color.rgb = RGBColor(244, 247, 250)
        
        img_path = os.path.join(images_path, data["image"])
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(7), Inches(1), width=Inches(5.5))

    filename = f"public/Apresentacao_Foundry_{module_name.replace(' ', '_')}.pptx"
    prs.save(filename)
    print(f"Sucesso: {filename}")

if __name__ == "__main__":
    if not os.path.exists("public"): os.makedirs("public")
    create_pptx("Modulo 01", M01)
    create_pptx("Modulo 02", M02)
    create_pptx("Modulo 03", M03)
    create_pptx("Modulo 04", M04)
