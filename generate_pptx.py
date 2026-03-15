from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Dados dos slides (extracao direta do slidesData.js)
slides_data = [
    {"id": 1, "title": "Azure AI Foundry", "subtitle": "A Nova Era dos Agentes Inteligentes", "description": "Seja bem-vindo ao curso definitivo sobre o ecossistema de IA da Microsoft. Aqui, transformamos conceitos complexos em agentes de mercado prontos para o mundo real.", "image": "img14.png"},
    {"id": 2, "title": "O que são Agentes de IA?", "description": "Assistentes digitais que não apenas respondem, mas agem. Eles percebem o ambiente, processam informações e executam tarefas multi-etapas de forma autônoma.", "image": "img4.png"},
    {"id": 3, "title": "Impacto nas Indústrias", "description": "Da saúde ao varejo, os agentes estão otimizando recursos, traduzindo idiomas em tempo real e criando experiências personalizadas 24/7.", "image": "img15.png"},
    {"id": 4, "title": "O Loop Percepção-Ação", "description": "O ciclo vital: Receber entrada (Percepção), analisar o contexto (Processamento) e realizar a tarefa no ecossistema (Ação).", "image": "img17.png"},
    {"id": 5, "title": "Agentes Reativos: Velocidade", "description": "Focados em respostas rápidas e diretas a estímulos imediatos. Ideais para FAQs e automações simples de atendimento.", "image": "img20.png"},
    {"id": 6, "title": "Agentes Deliberativos: Razão", "description": "Capazes de raciocínio complexo, planejamento e análise de dados. Essenciais para consultoria financeira e diagnósticos técnicos.", "image": "img11.png"},
    {"id": 7, "title": "Agentes Híbridos: Equilíbrio", "description": "Combinam a resposta instantânea dos reativos com a profundidade dos deliberativos. O topo da inteligência artificial.", "image": "img14.png"},
    {"id": 8, "title": "Escolhendo o Agente Ideal", "description": "Use reativos para rapidez, deliberativos para profundidade e híbridos para tarefas que exigem ambos com empatia.", "image": "img1.png"},
    {"id": 9, "title": "Arquitetura Unificada", "description": "Azure AI Foundry oferece escalabilidade global e segurança robusta para todo o ciclo de vida da IA, do protótipo à produção.", "image": "img22.png"},
    {"id": 10, "title": "Ambiente: ai.azure.com", "description": "Seu centro de comando unificado. Gerencie modelos, dados e implantações em uma interface web poderosa.", "image": "img10.png"},
    {"id": 11, "title": "Model Catalog", "description": "Mais de 11.000 modelos: de GPT-4 e Llama a modelos especializados de parceiros como Hugging Face e Meta.", "image": "img16.png"},
    {"id": 12, "title": "Sandbox Playground", "description": "O laboratório para inovação rápida. Ajuste parâmetros e valide ideias sem precisar de um projeto completo ou código inicial.", "image": "img21.png"},
    {"id": 13, "title": "No-Code para Pro-Code", "description": "Transição suave do portal visual para o VS Code. Use SDKs e extensões para levar sua IA a um nível profissional.", "image": "img13.png"},
    {"id": 14, "title": "Segurança Zero-Trust", "description": "Arquitetura de segurança máxima. Login protegido via Entra ID e isolamento completo de redes e chaves de API.", "image": "img7.png"},
    {"id": 15, "title": "Privacidade de Dados", "description": "Seus dados são seus. No Azure, informações de clientes e prompts não são usados para treinar modelos base de terceiros.", "image": "img8.png"},
    {"id": 16, "title": "Inteligência Responsável", "description": "Governança integrada. Filtros de segurança e verificações de alucinação para garantir uma IA ética e confiável.", "image": "img12.png"},
    {"id": 17, "title": "Hubs e Projetos", "description": "Estrutura organizacional eficiente. Hubs para recursos compartilhados e Projetos para experimentos isolados e seguros.", "image": "/img18.png"},
    {"id": 18, "title": "RAG & Grounding", "description": "Grounding: Conecte seu modelo aos seus próprios dados de negócio usando Retrieval Augmented Generation para respostas reais.", "image": "img19.png"},
    {"id": 19, "title": "Azure AI Search", "description": "O motor de busca inteligente por trás dos agentes RAG. Rankeamento semântico e busca vetorial de alta performance.", "image": "img22.png"},
    {"id": 20, "title": "Agentic Retrieval", "description": "Evolução do RAG: o agente planeja múltiplos passos de busca para responder perguntas complexas com máxima precisão.", "image": "img3.png"},
    {"id": 21, "title": "Serviços de Linguagem", "description": "Tradução, resumo e extração de entidades em tempo real. Entenda e processe textos em escala global.", "image": "img5.png"},
    {"id": 22, "title": "Serviços de Visão", "description": "Permita que seus agentes 'enxerguem'. Extração de texto de formulários (OCR) e identificação de objetos em imagens.", "image": "img12.png"},
    {"id": 23, "title": "Serviços de Fala", "description": "Interação por voz. Síntese de voz natural e transcrição precisa em múltiplos sotaques e idiomas.", "image": "img4.png"},
    {"id": 24, "title": "Chat Playground", "description": "Explore o comportamento do modelo. Compare GPT-4 com outros modelos em tempo real com prompts personalizados.", "image": "img21.png"},
    {"id": 25, "title": "Tradução Multimodal", "description": "Combine tradução de áudio, texto e documentos mantendo a fidelidade e o contexto original da informação.", "image": "img6.png"},
    {"id": 26, "title": "DALL-E no Azure", "description": "Geração de imagens generativas via API. Crie artes técnicas e ilustrações sob demanda para suas aplicações.", "image": "img12.png"},
    {"id": 27, "title": "Gestão de Custos", "description": "Controle granular. Use os leaderboards e métricas de consumo para otimizar o custo-benefício dos seus agentes.", "image": "img10.png"},
    {"id": 28, "title": "MLOps em Escala", "description": "Gerencie o ciclo completo: desde o versionamento do modelo até o monitoramento de performance em produção.", "image": "img1.png"},
    {"id": 29, "title": "Deployments Robustos", "description": "Transforme experimentos em serviços reais com endpoints de baixa latência e alta disponibilidade.", "image": "img22.png"},
    {"id": 30, "title": "O Futuro da Carreira em IA", "description": "Azure AI Foundry é o acelerador definitivo. Lidere a inovação criando soluções que resolvem problemas reais hoje.", "image": "img9.png"}
]

def create_pptx():
    prs = Presentation()
    
    # Definir tamanho do slide 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    images_path = os.path.join("public", "images")

    for data in slides_data:
        # Layout vazio para controle total
        slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(slide_layout)
        
        # Fundo escuro (aproximacao do design web)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(11, 14, 20)
        
        # Adicionar Titulo
        left = Inches(0.5)
        top = Inches(0.5)
        width = Inches(6)
        height = Inches(1.5)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = data["title"]
        p.font.bold = True
        p.font.size = Pt(44)
        p.font.color.rgb = RGBColor(0, 120, 212) # Azul Azure
        
        # Adicionar Descricao
        top = Inches(2.2)
        height = Inches(4)
        txBoxDesc = slide.shapes.add_textbox(left, top, width, height)
        tfDesc = txBoxDesc.text_frame
        tfDesc.word_wrap = True
        pDesc = tfDesc.add_paragraph()
        pDesc.text = data["description"]
        pDesc.font.size = Pt(22)
        pDesc.font.color.rgb = RGBColor(244, 247, 250)
        
        # Adicionar Imagem
        img_name = data["image"].replace("/images/", "")
        img_path = os.path.join(images_path, img_name)
        
        if os.path.exists(img_path):
            img_left = Inches(7)
            img_top = Inches(1)
            img_width = Inches(5.5)
            # A altura sera proporcional se nao for definida
            slide.shapes.add_picture(img_path, img_left, img_top, width=img_width)
        else:
            print(f"Aviso: Imagem nao encontrada {img_path}")

    output_file = "Apresentacao_Azure_AI_Foundry.pptx"
    prs.save(output_file)
    print(f"Sucesso: {output_file} criado!")

if __name__ == "__main__":
    create_pptx()
